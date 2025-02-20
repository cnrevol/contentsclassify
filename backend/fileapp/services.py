import mimetypes
import os
from django.core.files.storage import default_storage
from django.core.files.base import ContentFile
import hashlib
import PyPDF2
import docx
import openpyxl
from PIL import Image
import pytesseract
from pptx import Presentation
from classifier.models import ClassificationResult
from classifier.services import ContentClassificationService
from .models import ProcessedFile, TextInput
from django.conf import settings

class FileProcessingService:
    def __init__(self):
        self.classification_service = ContentClassificationService()

    @staticmethod
    def process_file(file):
        # Generate a unique path for the file
        file_path = default_storage.save(f'uploads/{file.name}', ContentFile(file.read()))
        
        # Get the file type using mimetypes
        file_type, _ = mimetypes.guess_type(file.name)
        if not file_type:
            file_type = 'application/octet-stream'
        
        # Get file size
        file_size = file.size
        
        # Generate content hash
        file.seek(0)
        content_hash = hashlib.md5(file.read()).hexdigest()
        
        # Extract text content (placeholder)
        extracted_text = "Text extraction will be implemented later"
        
        # Process metadata
        metadata = {
            'size': file_size,
            'type': file_type,
            'name': file.name
        }
        
        return {
            'file_path': file_path,
            'file_type': file_type,
            'content_hash': content_hash,
            'extracted_text': extracted_text,
            'metadata': metadata
        }

    @staticmethod
    def delete_file(file_path):
        if default_storage.exists(file_path):
            default_storage.delete(file_path)

    def process_file(self, file_instance):
        """Process uploaded file"""
        try:
            # Detect file type
            file_type = self._detect_file_type(file_instance.file_name)
            file_instance.file_type = file_type
            
            # Extract text based on file type
            extracted_text = self._extract_text(file_instance.file, file_type)
            file_instance.extracted_text = extracted_text
            
            # Generate content hash
            content_hash = self._generate_hash(extracted_text)
            
            # Create classification result
            classification_result = self._classify_content(
                content_hash,
                extracted_text,
                file_type,
                file_instance.user
            )
            
            # Update file instance
            file_instance.classification_result = classification_result[0]
            file_instance.processed = True
            file_instance.save()
            
        except Exception as e:
            if not file_instance.metadata:
                file_instance.metadata = {}
            file_instance.metadata['processing_error'] = str(e)
            file_instance.save()
            raise
            
    def process_text(self, text_instance):
        """Process text input"""
        try:
            # Generate content hash
            content_hash = self._generate_hash(text_instance.content)
            
            # Create classification results
            classification_results = self._classify_content(
                content_hash,
                text_instance.content,
                'text',
                text_instance.user
            )
            
            # Get LLM info from the first result (they're all the same for one submission)
            if classification_results and len(classification_results) > 0:
                first_result = classification_results[0]
                text_instance.llm_provider = first_result.metadata.get('llm_provider', 'unknown')
                text_instance.llm_model = first_result.metadata.get('llm_model', 'unknown')
                # Remove LLM info from individual results
                for result in classification_results:
                    if 'llm_provider' in result.metadata:
                        del result.metadata['llm_provider']
                    if 'llm_model' in result.metadata:
                        del result.metadata['llm_model']
                    result.save()
            
            # Update text instance
            text_instance.processed = True
            text_instance.save()
            
            # Add all classification results
            if isinstance(classification_results, list):
                text_instance.classification_results.set(classification_results)
            elif classification_results:
                text_instance.classification_results.set([classification_results])
            
        except Exception as e:
            raise
            
    def _detect_file_type(self, filename):
        """Detect file type using mimetypes"""
        file_type, _ = mimetypes.guess_type(filename)
        if not file_type:
            # Get file extension
            _, ext = os.path.splitext(filename)
            ext = ext.lower()
            
            # Map common extensions to file types
            type_map = {
                '.txt': 'text',
                '.pdf': 'pdf',
                '.doc': 'doc',
                '.docx': 'doc',
                '.xls': 'xls',
                '.xlsx': 'xls',
                '.ppt': 'ppt',
                '.pptx': 'ppt',
                '.jpg': 'image',
                '.jpeg': 'image',
                '.png': 'image',
                '.gif': 'image',
            }
            return type_map.get(ext, 'other')
            
        # Map mime types to our file types
        mime_map = {
            'text/': 'text',
            'application/pdf': 'pdf',
            'application/msword': 'doc',
            'application/vnd.openxmlformats-officedocument.wordprocessingml': 'doc',
            'application/vnd.ms-excel': 'xls',
            'application/vnd.openxmlformats-officedocument.spreadsheetml': 'xls',
            'application/vnd.ms-powerpoint': 'ppt',
            'application/vnd.openxmlformats-officedocument.presentationml': 'ppt',
            'image/': 'image',
        }
        
        for mime_prefix, file_type in mime_map.items():
            if file_type.startswith(mime_prefix):
                return file_type
                
        return 'other'
        
    def _extract_text(self, file, file_type):
        """Extract text from file based on its type"""
        try:
            file_path = file.path
        except:
            # If file.path is not available, try to get the path from storage
            file_path = default_storage.path(file.name)
        
        try:
            if file_type == 'text':
                with open(file_path, 'r', encoding='utf-8') as f:
                    return f.read()
                    
            elif file_type == 'pdf':
                text = []
                with open(file_path, 'rb') as f:
                    pdf = PyPDF2.PdfReader(f)
                    for page in pdf.pages:
                        text.append(page.extract_text())
                return '\n'.join(text)
                
            elif file_type == 'doc':
                doc = docx.Document(file_path)
                return '\n'.join([paragraph.text for paragraph in doc.paragraphs])
                
            elif file_type == 'xls':
                text = []
                wb = openpyxl.load_workbook(file_path)
                for sheet in wb.sheetnames:
                    ws = wb[sheet]
                    for row in ws.rows:
                        text.extend([str(cell.value) for cell in row if cell.value])
                return '\n'.join(text)
                
            elif file_type == 'ppt':
                text = []
                prs = Presentation(file_path)
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text.append(shape.text)
                return '\n'.join(text)
                
            elif file_type == 'image':
                return pytesseract.image_to_string(Image.open(file_path))
                
            return ''
        except Exception as e:
            return f'Error extracting text: {str(e)}'
        
    def _generate_hash(self, content):
        """Generate SHA-256 hash of content"""
        if isinstance(content, str):
            content = content.encode()
        return hashlib.sha256(content).hexdigest()
        
    def _classify_content(self, content_hash, content, content_type, user):
        """Classify content using LLM"""
        # Get classification from LLM
        results = self.classification_service.classify_content(content, content_type)
        
        # Create classification results for each group
        classification_results = []
        for result in results:
            classification_result = ClassificationResult.objects.create(
                content_type=content_type,
                content_hash=content_hash,
                classification=result['classification'],
                confidence=result['confidence'],
                metadata={
                    'explanation': result['explanation'],
                    'category_group_name': result['category_group_name']
                },
                user=user,
                category_group_id=result.get('category_group_id'),
                llm_provider=result.get('llm_provider', 'unknown'),
                llm_model=result.get('llm_model', 'unknown')
            )
            classification_results.append(classification_result)
        
        return classification_results

    def process_content(self, content: str, content_type: str, user=None, llm_provider=None):
        """
        Process the content using the classification service
        
        Args:
            content (str): Content to process
            content_type (str): Type of content
            user (User, optional): User who initiated the request
            llm_provider (str, optional): LLM provider to use for classification
        """
        try:
            # Update LLM provider if specified
            if llm_provider:
                provider_config = settings.LLM_PROVIDERS.get(llm_provider)
                if not provider_config:
                    raise ValueError(f"Unsupported LLM provider: {llm_provider}")
                
                # Get provider-specific configuration
                config_prefix = provider_config.get('config_prefix', llm_provider.upper())
                config = {
                    'api_key': getattr(settings, f"{config_prefix}_API_KEY", ''),
                    'base_url': getattr(settings, f"{config_prefix}_API_URL", ''),
                    'model_name': provider_config['default_model'],
                    'temperature': settings.DEFAULT_LLM_CONFIG['temperature'],
                    'max_tokens': settings.DEFAULT_LLM_CONFIG['max_tokens']
                }
                
                self.classification_service.update_provider(llm_provider, config)
                
            # Classify content
            results = self.classification_service.classify_content(content, content_type)
            
            # Create classification results
            classification_results = []
            for result in results:
                classification_result = ClassificationResult.objects.create(
                    content_type=content_type,
                    content_hash=self._generate_hash(content),
                    classification=result['classification'],
                    confidence=result['confidence'],
                    metadata={
                        'explanation': result['explanation'],
                        'category_group_name': result['category_group_name']
                    },
                    user=user,
                    category_group_id=result.get('category_group_id'),
                    llm_provider=result.get('llm_provider', 'unknown'),
                    llm_model=result.get('llm_model', 'unknown')
                )
                classification_results.append(classification_result)
            
            return classification_results
            
        except Exception as e:
            raise 